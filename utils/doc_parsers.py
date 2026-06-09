# utils/doc_parsers.py
import os
from io import BytesIO


def _decrypt_office_file(file_path: str, password: str | None) -> BytesIO | None:
    """
    Try to decrypt OOXML / supported Office files with msoffcrypto.
    Return a BytesIO object when the file is encrypted and successfully decrypted.
    Return None when the file is not encrypted or msoffcrypto is unavailable.
    """
    try:
        import msoffcrypto
    except ImportError:
        return None

    try:
        with open(file_path, "rb") as f:
            office_file = msoffcrypto.OfficeFile(f)
            if not office_file.is_encrypted():
                return None
            if not password:
                raise ValueError("Office文档受加密保护，需要密码")
            office_file.load_key(password=password)
            decrypted = BytesIO()
            office_file.decrypt(decrypted)
            decrypted.seek(0)
            return decrypted
    except ValueError:
        raise
    except Exception as e:
        # Some normal, unencrypted Office files may not be recognized by msoffcrypto.
        # Only surface the error when a password was provided or encryption was likely.
        if password:
            raise ValueError(f"Office文档解密失败: {e}")
        return None


def parse_txt(file_path: str):
    """
    解析 TXT 文件。
    返回生成器，每次产出（行号，文本）
    """
    encodings = ['utf-8', 'gb18030', 'gbk', 'gb2312', 'utf-16']
    for enc in encodings:
        try:
            with open(file_path, 'r', encoding=enc) as f:
                for line_num, line in enumerate(f, start=1):
                    text = line.strip()
                    if text:
                        yield line_num, text
            return
        except UnicodeDecodeError:
            continue

    raise ValueError("无法识别文件编码或文件已损坏")


def parse_docx(file_path: str, password: str | None = None):
    """
    解析 Word(.docx) 文件。
    若文件加密且安装 msoffcrypto-tool，则可传入 password 解密后继续读取。
    Word 非严格“行”换行，按 Paragraph 作为行号基准。
    """
    try:
        import docx

        decrypted = _decrypt_office_file(file_path, password)
        doc = docx.Document(decrypted if decrypted is not None else file_path)
        for i, para in enumerate(doc.paragraphs, start=1):
            text = para.text.strip()
            if text:
                yield i, text
    except ImportError as e:
        raise ValueError(f"缺少 DOCX 解析依赖: {e}")
    except Exception as e:
        raise ValueError(f"文档受保护(加密)、格式已损坏或读取失败: {e}")


def parse_xlsx(file_path: str, password: str | None = None):
    """
    解析 Excel(.xlsx) 文件。
    返回：（表名+行号，该行所有单元格拼接的文本）。
    若文件加密且安装 msoffcrypto-tool，则可传入 password 解密后继续读取。
    """
    try:
        import openpyxl

        decrypted = _decrypt_office_file(file_path, password)
        source = decrypted if decrypted is not None else file_path
        wb = openpyxl.load_workbook(source, read_only=True, data_only=True)
        try:
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                for row_idx, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                    row_texts = [str(cell) for cell in row if cell is not None and str(cell).strip()]
                    if row_texts:
                        text = " | ".join(row_texts)
                        yield f"表[{sheet_name}] - 第{row_idx}行", text
        finally:
            wb.close()
    except ImportError as e:
        raise ValueError(f"缺少 XLSX 解析依赖: {e}")
    except Exception as e:
        raise ValueError(f"Excel读取失败或受保护: {str(e)}")


def parse_pptx(file_path: str, password: str | None = None):
    """
    解析 PowerPoint (.pptx) 文件。
    返回: (页码, 该页所有文本框拼接的文本)。
    若文件加密且安装 msoffcrypto-tool，则可传入 password 解密后继续读取。
    """
    try:
        import pptx

        decrypted = _decrypt_office_file(file_path, password)
        prs = pptx.Presentation(decrypted if decrypted is not None else file_path)
        for i, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)

            if slide_text:
                text = " ".join(slide_text).replace('\n', ' ')
                if text.strip():
                    yield f"第{i}页", text
    except ImportError as e:
        raise ValueError(f"缺少 PPTX 解析依赖: {e}")
    except Exception as e:
        raise ValueError(f"PPT读取失败或受保护: {e}")


def parse_pdf(file_path: str, password: str | None = None):
    """
    解析 PDF (.pdf) 文件。
    返回: (页码, 该页提取的所有文本)。
    若 PDF 加密，则支持传入 password 后继续解析。
    """
    doc = None
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(file_path)
        if doc.needs_pass:
            if not password:
                raise ValueError("PDF受加密保护，需要密码")
            if not doc.authenticate(password):
                raise ValueError("PDF密码错误或无法解密")

        for i, page in enumerate(doc, start=1):
            text = page.get_text("text").strip()
            if text:
                text = text.replace('\n', ' ')
                yield f"第{i}页", text
    except ImportError as e:
        raise ValueError(f"缺少 PDF 解析依赖: {e}")
    except Exception as e:
        raise ValueError(f"PDF读取失败或受加密保护: {e}")
    finally:
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass


def parse_xls(file_path: str, password: str | None = None):
    """解析老版本 Excel (.xls) 文件；若 msoffcrypto 可解密，则支持 password。"""
    try:
        import xlrd

        decrypted = _decrypt_office_file(file_path, password)
        if decrypted is not None:
            wb = xlrd.open_workbook(file_contents=decrypted.getvalue())
        else:
            wb = xlrd.open_workbook(file_path)
        for sheet in wb.sheets():
            for row_idx in range(sheet.nrows):
                row_values = sheet.row_values(row_idx)
                row_texts = [str(cell).strip() for cell in row_values if str(cell).strip()]
                if row_texts:
                    text = " | ".join(row_texts)
                    yield f"表[{sheet.name}] 第{row_idx+1}行", text
    except ImportError as e:
        raise ValueError(f"缺少 XLS 解析依赖: {e}")
    except Exception as e:
        raise ValueError(f"老版本Excel读取失败(可能加密或损坏): {e}")


def parse_doc(file_path: str, password: str | None = None):
    """使用 Windows COM 接口解析老版本 Word (.doc) 文件；尽量支持 PasswordDocument。"""
    word_app = None
    doc = None
    pythoncom = None
    try:
        import pythoncom
        import win32com.client

        pythoncom.CoInitialize()
        abs_path = os.path.abspath(file_path)

        word_app = win32com.client.DispatchEx("Word.Application")
        word_app.Visible = False
        word_app.DisplayAlerts = 0

        open_kwargs = {
            "ConfirmConversions": False,
            "ReadOnly": True,
            "Visible": False,
        }
        if password:
            open_kwargs["PasswordDocument"] = password

        doc = word_app.Documents.Open(abs_path, **open_kwargs)

        for i, para in enumerate(doc.Paragraphs, start=1):
            text = para.Range.Text.strip()
            if text:
                yield i, text
    except ImportError as e:
        raise ValueError(f"缺少旧版 Word 解析依赖: {e}")
    except Exception as e:
        raise ValueError(f"老版本Word读取失败(格式不匹配、加密或损坏): {e}")
    finally:
        if doc:
            try:
                doc.Close(False)
            except Exception:
                pass
        if word_app:
            try:
                word_app.Quit()
            except Exception:
                pass
        if pythoncom:
            pythoncom.CoUninitialize()


def parse_ppt(file_path: str, password: str | None = None):
    """使用 Windows COM 接口解析老版本 PowerPoint (.ppt) 文件；旧版 PPT 密码支持随 Office 版本而异。"""
    ppt_app = None
    presentation = None
    pythoncom = None
    try:
        import pythoncom
        import win32com.client

        pythoncom.CoInitialize()
        abs_path = os.path.abspath(file_path)

        ppt_app = win32com.client.DispatchEx("PowerPoint.Application")
        if password:
            try:
                presentation = ppt_app.Presentations.Open(abs_path, ReadOnly=True, WithWindow=False, Password=password)
            except TypeError:
                presentation = ppt_app.Presentations.Open(abs_path, ReadOnly=True, WithWindow=False)
        else:
            presentation = ppt_app.Presentations.Open(abs_path, ReadOnly=True, WithWindow=False)

        for i, slide in enumerate(presentation.Slides, start=1):
            slide_text = []
            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    slide_text.append(shape.TextFrame.TextRange.Text.strip())

            if slide_text:
                text = " ".join(slide_text).replace('\n', ' ').replace('\r', ' ')
                if text.strip():
                    yield f"第{i}页", text
    except ImportError as e:
        raise ValueError(f"缺少旧版 PPT 解析依赖: {e}")
    except Exception as e:
        raise ValueError(f"老版本PPT读取失败(可能加密或损坏): {e}")
    finally:
        if presentation:
            try:
                presentation.Close()
            except Exception:
                pass
        if ppt_app:
            try:
                ppt_app.Quit()
            except Exception:
                pass
        if pythoncom:
            pythoncom.CoUninitialize()
